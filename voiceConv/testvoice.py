import collections 
import keyboard
import collections.abc
from pptx import Presentation
import pyttsx3
from os import startfile
from pyautogui import click
from keyboard import press
from keyboard import write
from time import sleep
import datetime
from time import sleep
from pygame import mixer
from pyautogui import hotkey
import os
from datetime import *
from textToHindi import translatorText

files = os.listdir('.')
print(files)
for file in files:
    if file.endswith(".pptx"):
        filedir=file
        break
print(filedir)
def playmusic(file_path):
    mixer.init()
    # Loading the song
    mixer.music.load(file_path)  
    # Setting the volume
    mixer.music.set_volume(0.2)
    # Start playing the song
    mixer.music.play()
    while True:
        sleep(2)
        break

# startfile("C:\\Users\\Vinay\\AppData\\Local\\WhatsApp\\WhatsApp.exe")
startfile(filedir)
sleep(6)
click(x=1310, y=131)
sleep(2)
press('esc')
press('f5')

# file = 'ttt.pptx'
# file="C:\\Users\\Vinay Mishra\\Desktop\\PythonTraining\\VVNT_Tool_Github.pptx"

ppt=Presentation(filedir)

notes = []

for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note. 
    textNote = slide.notes_slide.notes_text_frame.text
    notes.append((page,textNote)) 

# initializing or creating object
engine = pyttsx3.init("sapi5")
# voices will contain voices like in this system there is 3 type of voice
voices = engine.getProperty("voices")

print(f"voices id is {voices[-1].id} ")
# print(voices)
# here we are choosing a voice or setting voice
engine.setProperty("voice", voices[1].id)

engine.setProperty('rate', 165)
# text = 'Your Text is here and this is bla bla bla bla'
# engine.save_to_file(text, 'name.mp3')
engine.runAndWait() # don't forget to use this line
# engine.setProperty('volume', 2.0)
# engine.save_to_file('Hi i am vinay and i work in vvnt as automation engineer', 'audio.mp3')
def wish_me():
    speak("Hello Everyone.")

def speak(audio):
    engine.say(audio)
    engine.runAndWait()
# print(notes)
keyboard.press('alt+f7')
sleep(1)
keyboard.release('alt+f7')
today = str(date.today())
now = datetime.now()
current_time = str(now.strftime("%H-%M"))
sleep(4)
wish_me()
for i,j in notes:
    # print(f"index is {i} and the data is {j}")
    # speak(j) editing here for adding new code
    
    # print(j)
    # sleep(200)    w2q
    if i==(len(notes)-1):
        sleep(2)
        break
    else:
        translatorText(j)
    # sleep(1)
    press('right')
    playmusic("success.mp3")
    # sleep(1)
keyboard.press('alt+f8')
sleep(1)
keyboard.release('alt+f8')
sleep(1)
keyboard.press('alt+f4')
keyboard.release('alt+f4')
sleep(1)
keyboard.press('alt+f4')
keyboard.release('alt+f4')
sleep(1)
# opening downloaded file
files2 = os.listdir('C:\\Users\\Administrator\\Videos')
print(f"files in videos are {files2}")
filename=f"{today} {current_time}"
for i in files2:
    if filename in str(i):
        os.startfile(f"C:\\Users\\Administrator\\Videos\\{i}")
        break
sleep(1)
press('f')